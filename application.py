from flask import Flask,render_template,request,send_file,url_for
from google_images_download import google_images_download
from pptx import Presentation
from pptx.util import Inches


import sys

application=Flask(__name__)
app=application

@app.route('/')

def homepage():
	return render_template("index_new.html")

@app.route('/ppt')

def genppt():
	return render_template("ppt_free.html")

@app.route('/contact')

def contactpage():
	return render_template("contacts.html")

@app.route('/generate',methods=['POST','GET'])

def generate_ppt():
	paths=[]
	result=request.args.get('input')
	arguments={"keywords": result,"limit":5}
	response=google_images_download.googleimagesdownload()
	paths.append(response.download(arguments))
	img1= paths[0][arguments['keywords']][0]
	prs=Presentation()
	slide_layout1=prs.slide_layouts[0]
	slide1=prs.slides.add_slide(slide_layout1)
	shapes=slide1.shapes
	left_pic=Inches(0)
	top_pic=Inches(0)
	width_pic=Inches(3)
	height_pic=Inches(7.5)
	pic=slide1.shapes.add_picture(img1,left_pic,top_pic,width_pic,height_pic)
	placeholder_title=slide1.placeholders[0]
	placeholder_title.left=Inches(3.3)
	placeholder_title.top=Inches(0.5)
	placeholder_title.width=Inches(6)
	placeholder_title.height=Inches(2)
	placeholder_content=slide1.placeholders[1]
	placeholder_content.left=Inches(3.4)
	placeholder_content.top=Inches(4)
	placeholder_content.width=Inches(6)
	placeholder_content.height=Inches(2.3)
	slide_layout2=prs.slide_layouts[0]
	slide2=prs.slides.add_slide(slide_layout1)
	shapes=slide2.shapes
	left_pic=Inches(0)
	top_pic=Inches(0)
	width_pic=Inches(3)
	height_pic=Inches(7.5)
	img2= paths[0][arguments['keywords']][1]
	pic=slide2.shapes.add_picture(img2,left_pic,top_pic,width_pic,height_pic)
	placeholder_title=slide2.placeholders[0]
	placeholder_title.left=Inches(3.3)
	placeholder_title.top=Inches(0.5)
	placeholder_title.width=Inches(6)
	placeholder_title.height=Inches(2)
	placeholder_content=slide2.placeholders[1]
	placeholder_content.left=Inches(3.4)
	placeholder_content.top=Inches(4)
	placeholder_content.width=Inches(6)
	placeholder_content.height=Inches(2.3)
	slide_layout3=prs.slide_layouts[0]
	slide3=prs.slides.add_slide(slide_layout1)
	shapes=slide3.shapes
	left_pic=Inches(0)
	top_pic=Inches(0)
	width_pic=Inches(3)
	height_pic=Inches(7.5)
	img3= paths[0][arguments['keywords']][2]
	pic=slide3.shapes.add_picture(img3,left_pic,top_pic,width_pic,height_pic)
	placeholder_title=slide3.placeholders[0]
	placeholder_title.left=Inches(3.3)
	placeholder_title.top=Inches(0.5)
	placeholder_title.width=Inches(6)
	placeholder_title.height=Inches(2)
	placeholder_content=slide3.placeholders[1]
	placeholder_content.left=Inches(3.4)
	placeholder_content.top=Inches(4)
	placeholder_content.width=Inches(6)
	placeholder_content.height=Inches(2.3)
	slide_layout4=prs.slide_layouts[0]
	slide4=prs.slides.add_slide(slide_layout1)
	shapes=slide4.shapes
	left_pic=Inches(0)
	top_pic=Inches(0)
	width_pic=Inches(3)
	height_pic=Inches(7.5)
	img4= paths[0][arguments['keywords']][3]
	pic=slide4.shapes.add_picture(img4,left_pic,top_pic,width_pic,height_pic)
	placeholder_title=slide4.placeholders[0]
	placeholder_title.left=Inches(3.3)
	placeholder_title.top=Inches(0.5)
	placeholder_title.width=Inches(6)
	placeholder_title.height=Inches(2)
	placeholder_content=slide4.placeholders[1]
	placeholder_content.left=Inches(3.4)
	placeholder_content.top=Inches(4)
	placeholder_content.width=Inches(6)
	placeholder_content.height=Inches(2.3)
	slide_layout5=prs.slide_layouts[0]
	slide5=prs.slides.add_slide(slide_layout1)
	shapes=slide5.shapes
	left_pic=Inches(0)
	top_pic=Inches(0)
	width_pic=Inches(3)
	height_pic=Inches(7.5)
	img5= paths[0][arguments['keywords']][4]
	pic=slide5.shapes.add_picture(img5,left_pic,top_pic,width_pic,height_pic)
	placeholder_title=slide5.placeholders[0]
	placeholder_title.left=Inches(3.3)
	placeholder_title.top=Inches(0.5)
	placeholder_title.width=Inches(6)
	placeholder_title.height=Inches(2)
	placeholder_content=slide5.placeholders[1]
	placeholder_content.left=Inches(3.4)
	placeholder_content.top=Inches(4)
	placeholder_content.width=Inches(6)
	placeholder_content.height=Inches(2.3)
  	
    
	prs.save("result.pptx")
	print("file saved")
    
	return send_file('result.pptx',as_attachment=True)



if __name__=="__main__":
	app.run(debug=True)


